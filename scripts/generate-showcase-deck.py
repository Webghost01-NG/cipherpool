#!/usr/bin/env python3
"""Generate the editable Veylott showcase deck from verified project facts."""

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "docs" / "showcase" / "presentation"
SCREENSHOT = OUTPUT_DIR / "assets" / "live-dashboard.png"
OUTPUT = OUTPUT_DIR / "Veylott-Presentation.pptx"

WHITE = RGBColor(255, 255, 255)
INK = RGBColor(20, 31, 55)
MUTED = RGBColor(93, 110, 142)
BLUE = RGBColor(49, 87, 246)
BLUE_DARK = RGBColor(24, 53, 159)
BLUE_PALE = RGBColor(239, 243, 255)
LINE = RGBColor(207, 217, 241)
GREEN = RGBColor(8, 127, 91)
GREEN_PALE = RGBColor(229, 249, 241)
AMBER = RGBColor(145, 82, 0)
AMBER_PALE = RGBColor(255, 244, 214)
FONT = "Liberation Sans"
MONO = "Liberation Mono"


def add_text(slide, text, x, y, w, h, *, size=18, color=INK, bold=False,
             font=FONT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0.02, line_spacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rich_text(slide, segments, x, y, w, h, *, size=18, color=INK,
                  align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, segment_color, bold in segments:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = segment_color or color
    return box


def add_link(slide, label, url, x, y, w, h, *, size=12, color=BLUE, font=MONO):
    box = add_text(slide, "", x, y, w, h, size=size, color=color, font=font)
    paragraph = box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = label
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    run.hyperlink.address = url
    return box


def add_box(slide, x, y, w, h, *, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1)
    return box


def add_pill(slide, text, x, y, w, *, fill=BLUE_PALE, color=BLUE):
    pill = add_box(slide, x, y, w, 0.32, fill=fill, line=fill)
    add_text(slide, text.upper(), x + 0.08, y + 0.02, w - 0.16, 0.24,
             size=8, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE)
    return pill


def add_brand(slide, number):
    mark = add_box(slide, 0.58, 0.28, 0.34, 0.34, fill=BLUE, line=BLUE)
    add_text(slide, "V", 0.58, 0.285, 0.34, 0.28, size=13, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Veylott", 1.02, 0.3, 1.5, 0.28, size=12, bold=True)
    add_text(slide, f"RESEARCH DEPLOYMENT · ETHEREUM SEPOLIA", 8.55, 0.34, 3.65, 0.2,
             size=7, color=MUTED, font=MONO, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return mark


if __name__ == "__main__":
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts.TITLE_SLIDE)
    mark = add_brand(slide, 1)
    prs.save(OUTPUT)
    print(f"Deck saved to: {OUTPUT.resolve()}")